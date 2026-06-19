"""Branded PowerPoint renderer (python-pptx). Sections become slides."""

from io import BytesIO

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from crm.documents.domain.enums import DOCUMENT_TYPE_LABELS
from crm.documents.domain.payload import format_money

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _rgb(hex_color: str, fallback: str) -> RGBColor:
    h = (hex_color or "").lstrip("#")
    try:
        return RGBColor.from_string(h if len(h) == 6 else fallback)
    except ValueError:
        return RGBColor.from_string(fallback)


def _blank(prs: Presentation):
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    return prs.slides.add_slide(layout)


def _text(
    slide, left, top, width, height, text, *, size=18, bold=False, color=None, align=PP_ALIGN.LEFT
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate((text or "").split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        if color is not None:
            run.font.color.rgb = color
    return box


def _bar(slide, color, *, height, top=0):
    shape = slide.shapes.add_shape(1, 0, top, SLIDE_W, height)  # 1 = rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def render_pptx(brand: dict, doc_type: str, payload: dict) -> bytes:
    label = DOCUMENT_TYPE_LABELS.get(doc_type, "Documento")
    primary = _rgb(brand.get("primary_color", "#7C6CFF"), "7C6CFF")
    accent = _rgb(brand.get("accent_color", "#22C55E"), "22C55E")
    white = RGBColor.from_string("FFFFFF")
    ink = _rgb(brand.get("text_color", "#16161D"), "16161D")
    currency = payload["currency"] or brand.get("currency", "")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # --- Title slide (full brand color) ---
    cover = _blank(prs)
    cover.background.fill.solid()
    cover.background.fill.fore_color.rgb = primary
    _text(
        cover,
        Inches(0.9),
        Inches(0.7),
        Inches(11),
        Inches(0.6),
        brand.get("business_name") or "",
        size=18,
        bold=True,
        color=white,
    )
    _text(
        cover,
        Inches(0.9),
        Inches(2.6),
        Inches(11.5),
        Inches(1.8),
        payload["title"] or label,
        size=40,
        bold=True,
        color=white,
    )
    if payload["subtitle"]:
        _text(
            cover,
            Inches(0.9),
            Inches(4.2),
            Inches(11.5),
            Inches(0.8),
            payload["subtitle"],
            size=20,
            color=white,
        )
    footer_bits = " · ".join(
        b for b in [label, payload["client"]["name"], payload["valid_until"]] if b
    )
    _text(
        cover,
        Inches(0.9),
        Inches(6.4),
        Inches(11.5),
        Inches(0.6),
        footer_bits,
        size=12,
        color=white,
    )

    # --- Section slides ---
    for section in payload["sections"]:
        slide = _blank(prs)
        _bar(slide, accent, height=Inches(0.14))
        _text(
            slide,
            Inches(0.9),
            Inches(0.6),
            Inches(11.5),
            Inches(1.0),
            section["heading"] or "",
            size=28,
            bold=True,
            color=primary,
        )
        _text(
            slide,
            Inches(0.9),
            Inches(1.8),
            Inches(11.5),
            Inches(5.0),
            section["body"] or "",
            size=18,
            color=ink,
        )

    # --- Items + totals slide ---
    if payload["has_items"]:
        slide = _blank(prs)
        _bar(slide, accent, height=Inches(0.14))
        _text(
            slide,
            Inches(0.9),
            Inches(0.5),
            Inches(11),
            Inches(0.8),
            "Detalle e inversión",
            size=26,
            bold=True,
            color=primary,
        )

        items = payload["items"]
        rows = len(items) + 1
        table_shape = slide.shapes.add_table(
            rows, 4, Inches(0.9), Inches(1.7), Inches(11.5), Inches(0.4 * rows)
        )
        table = table_shape.table
        table.columns[0].width = Inches(6.5)
        table.columns[1].width = Inches(1.4)
        table.columns[2].width = Inches(1.8)
        table.columns[3].width = Inches(1.8)
        headers = ["Detalle", "Cant.", "Precio unit.", "Importe"]
        for c, htext in enumerate(headers):
            cell = table.cell(0, c)
            cell.text = htext
            p = cell.text_frame.paragraphs[0]
            p.runs[0].font.size = Pt(12)
            p.runs[0].font.bold = True
            p.runs[0].font.color.rgb = white
            cell.fill.solid()
            cell.fill.fore_color.rgb = primary
        for r, it in enumerate(items, start=1):
            values = [
                it["description"] or "—",
                f"{it['quantity'].normalize():f}",
                format_money(it["unit_price"], currency),
                format_money(it["line_total"], currency),
            ]
            for c, val in enumerate(values):
                cell = table.cell(r, c)
                cell.text = val
                run = cell.text_frame.paragraphs[0].runs[0]
                run.font.size = Pt(11)
                run.font.color.rgb = ink

        total_line = f"Total: {format_money(payload['total'], currency)}"
        _text(
            slide,
            Inches(0.9),
            Inches(1.9) + Inches(0.4 * rows),
            Inches(11.5),
            Inches(0.7),
            total_line,
            size=18,
            bold=True,
            color=accent,
            align=PP_ALIGN.RIGHT,
        )

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
